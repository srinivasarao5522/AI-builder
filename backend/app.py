import os
import io
import json
import base64
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, Response, StreamingResponse
from pydantic import BaseModel
from dotenv import load_dotenv
import openai
import chromadb
from PyPDF2 import PdfReader
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation

# Load environment variables
load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
openai.api_key = OPENAI_API_KEY

app = FastAPI(title="AI CV Builder Backend")

# CORS setup
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize Vector DB (Chroma)
try:
    chroma_client = chromadb.Client()
    collection = chroma_client.create_collection(name="cv_templates")
    # Seed Vector DB with some generic beautiful templates
    templates_data = [
        {"id": "t1", "name": "Classic Elegance", "desc": "A traditional, clean professional layout perfect for corporate, finance, and legal roles. Emphasizes text hierarchy."},
        {"id": "t2", "name": "Modern Minimalist", "desc": "Sleek and contemporary, great for tech, engineering, and startup applications. Minimalist design with strong typography."},
        {"id": "t3", "name": "Creative Studio", "desc": "Colorful, vibrant layout with a left sidebar, ideal for designers, marketers, and creatives. Supports robust portfolio links."},
        {"id": "t4", "name": "Executive Standard", "desc": "Symmetrical, balanced and formal layout suitable for C-level executives. Features impactful achievement bullets."}
    ]
    collection.add(
        documents=[t["desc"] for t in templates_data],
        metadatas=[{"name": t["name"]} for t in templates_data],
        ids=[t["id"] for t in templates_data]
    )
except Exception as e:
    print(f"Warning: Chroma initialization failed: {e}")
    templates_data = []

class ChatRequest(BaseModel):
    message: str
    history: list = []
    cv_data: dict = {}
    language: str = "en"

@app.post("/api/chat")
async def chat_interaction(req: ChatRequest):
    if not openai.api_key:
        return {"reply": "API Key missing in backend. Please set OPENAI_API_KEY in .env.", "cv_data": req.cv_data}

    system_prompt = f"""You are an advanced multilingual CV/Resume builder AI assistant navigating a user through resume creation.
Language: {req.language}
Your goal is to guide the user to collect their info (Name, Contact, Summary, Experience, Education, Skills).
If they provide some, ask for what's logically next.
Keep the conversation incredibly friendly, helpful, encouraging and concise.

You must output in JSON format with exactly 2 keys:
"reply": Your conversational text back to the user (in the requested '{req.language}' language).
"cv_data": The fully accumulated and updated CV data structure (JSON object) reflecting EVERYTHING extracted so far. 

IMPORTANT: Ensure cv_data has the following standard keys: name, email, phone, summary, experience (list of objects with title, company, dates, description), education (list of objects), skills (list of strings).
Current CV Data: {json.dumps(req.cv_data)}
"""

    messages = [{"role": "system", "content": system_prompt}]
    for msg in req.history:
        messages.append(msg)
    messages.append({"role": "user", "content": req.message})

    try:
        response = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            response_format={ "type": "json_object" },
            temperature=0.7
        )
        content = response.choices[0].message.content
        parsed = json.loads(content)
        
        reply = parsed.get("reply", "I processed that. What's next?")
        new_data = parsed.get("cv_data", req.cv_data)

        # Template logic: if enough data but no template selected, we could trigger a suggestion
        return {"reply": reply, "cv_data": new_data}
    except Exception as e:
        return {"reply": f"Error communicating with AI: {str(e)}", "cv_data": req.cv_data}

class EnhanceRequest(BaseModel):
    text: str
    context: str = ""
    language: str = "en"

@app.post("/api/enhance")
async def enhance_text(req: EnhanceRequest):
    if not openai.api_key:
        return {"enhanced": req.text}
    try:
        sys_prompt = f"You are a professional resume writer. Your job is to take the user's short or poorly written text and auto-complete/enhance it into a punchy, professional, results-oriented resume bullet point. Expand upon their idea naturally. Output ONLY the beautifully rewritten text. Language: {req.language}"
        user_msg = f"Role Context: {req.context}\nInput Text: {req.text}"
        response = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "system", "content": sys_prompt}, {"role": "user", "content": user_msg}],
            temperature=0.7
        )
        return {"enhanced": response.choices[0].message.content.strip()}
    except Exception as e:
        return {"enhanced": req.text}

@app.post("/api/transcribe")
async def transcribe_audio(audio: UploadFile = File(...)):
    if not openai.api_key:
        return {"text": "API Key missing."}
    
    try:
        # Load audio into memory
        file_bytes = await audio.read()
        
        # Save temp file as whisper requires filename
        temp_filename = f"temp_{audio.filename}"
        with open(temp_filename, "wb") as f:
            f.write(file_bytes)
            
        with open(temp_filename, "rb") as audio_file:
            transcript = openai.audio.transcriptions.create(
                model="whisper-1", 
                file=audio_file
            )
            
        os.remove(temp_filename)
        return {"text": transcript.text}
    except Exception as e:
        return {"text": f"[Audio transcription failed: {str(e)}]"}

@app.post("/api/upload")
async def upload_document(file: UploadFile = File(...)):
    text = ""
    try:
        if file.filename.endswith('.pdf'):
            reader = PdfReader(file.file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        elif file.filename.endswith('.docx'):
            doc = Document(file.file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file.filename.endswith('.pptx'):
            ppt = Presentation(file.file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            file_bytes = await file.read()
            text = file_bytes.decode('utf-8', errors='ignore')
            
        # Optional: Ask GPT to pre-parse this text into cv_data immediately
        if openai.api_key and len(text.strip()) > 0:
            prompt = f"Parse the following resume text into a standard JSON schema containing keys: name, email, phone, summary, experience (list), education (list), skills (list).\n\nText:\n{text[:5000]}"
            resp = openai.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                response_format={"type": "json_object"}
            )
            parsed_data = json.loads(resp.choices[0].message.content)
            return {"extracted_text": text, "parsed_cv": parsed_data}

    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to parse document: {str(e)}")

    return {"extracted_text": text, "parsed_cv": {}}

@app.get("/api/templates")
async def get_templates(query: str = ""):
    if query:
        try:
            results = collection.query(query_texts=[query], n_results=2)
            matches = []
            for i in range(len(results["ids"][0])):
                matches.append({
                    "id": results["ids"][0][i],
                    "name": results["metadatas"][0][i]["name"],
                    "desc": results["documents"][0][i]
                })
            return {"templates": matches}
        except:
            pass
    return {"templates": templates_data}

@app.post("/api/generate")
async def generate_cv(req: dict):
    cv_data = req.get("cv_data", {})
    format_type = req.get("format", "json")
    
    if format_type == "json":
        return JSONResponse(content={"data": cv_data})

    if format_type == "docx":
        doc = Document()
        template_id = req.get("template_id", "t1")
        
        # Configure styles based on template
        styles = doc.styles
        heading1 = styles['Heading 1']
        heading2 = styles['Heading 2']
        normal = styles['Normal']
        
        # Base colors and alignment based on template
        main_color = RGBColor(0, 0, 0)
        accent_color = RGBColor(0, 0, 0)
        align_center = False
        
        if template_id == 't1': # Classic
            main_color = RGBColor(17, 17, 17) # Dark Gray
            heading1.font.name = 'Times New Roman'
            heading2.font.name = 'Times New Roman'
            normal.font.name = 'Times New Roman'
        elif template_id == 't2': # Modern
            main_color = RGBColor(15, 23, 42)
            accent_color = RGBColor(71, 85, 105)
            heading1.font.name = 'Arial'
            heading2.font.name = 'Arial'
            normal.font.name = 'Arial'
        elif template_id == 't3': # Creative
            main_color = RGBColor(15, 23, 42)
            accent_color = RGBColor(59, 130, 246) # Blue
            heading1.font.name = 'Calibri'
            heading2.font.name = 'Calibri'
            normal.font.name = 'Calibri'
        elif template_id == 't4': # Exec
            main_color = RGBColor(17, 17, 17)
            align_center = True
            heading1.font.name = 'Georgia'
            heading2.font.name = 'Georgia'
            normal.font.name = 'Georgia'
        
        heading1.font.color.rgb = main_color
        heading2.font.color.rgb = main_color

        # Inject Logo if present
        import urllib.request
        logo_data = req.get("logo_data", "")
        img_stream = None
        if logo_data:
            try:
                if logo_data.startswith("http"):
                    req_obj = urllib.request.Request(logo_data, headers={'User-Agent': 'Mozilla/5.0'})
                    with urllib.request.urlopen(req_obj) as response:
                        img_bytes = response.read()
                        img_stream = io.BytesIO(img_bytes)
                elif "base64," in logo_data:
                    img_bytes = base64.b64decode(logo_data.split(",")[1])
                    img_stream = io.BytesIO(img_bytes)
            except Exception as e:
                print(f"Failed to load logo: {e}")

        # Document Generator helper
        def add_summary(container):
            if cv_data.get("summary"):
                h = container.add_paragraph("Summary", style='Heading 1')
                if align_center: h.alignment = WD_ALIGN_PARAGRAPH.CENTER
                container.add_paragraph(cv_data.get("summary", ""))

        def add_exp(container):
            if cv_data.get("experience"):
                h = container.add_paragraph("Experience", style='Heading 1')
                if align_center: h.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for exp in cv_data.get("experience", []):
                    h2 = container.add_paragraph(exp.get("title", "Role"), style='Heading 2')
                    p = container.add_paragraph()
                    run = p.add_run(exp.get("company", "Company"))
                    run.bold = True
                    run.font.color.rgb = accent_color
                    if exp.get("dates"): p.add_run(f" | {exp.get('dates')}")
                    if exp.get("description"): container.add_paragraph(exp.get("description", ""))

        def add_edu(container):
            if cv_data.get("education"):
                h = container.add_paragraph("Education", style='Heading 1')
                if align_center: h.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for edu in cv_data.get("education", []):
                    container.add_paragraph(edu.get("degree", "Degree"), style='Heading 2')
                    container.add_paragraph(edu.get("institution", "Institution"))

        def add_skills(container):
            if cv_data.get("skills"):
                h = container.add_paragraph("Skills", style='Heading 1')
                if align_center: h.alignment = WD_ALIGN_PARAGRAPH.CENTER
                skills = cv_data.get("skills", [])
                p = container.add_paragraph(", ".join(skills) if isinstance(skills, list) else str(skills))
                if align_center: p.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Layout logic
        if template_id in ['t2', 't3']: # Side by side 2 column layouts
            table = doc.add_table(rows=1, cols=2)
            table.autofit = True
            left_cell = table.cell(0, 0)
            right_cell = table.cell(0, 1)

            if img_stream:
                p = left_cell.paragraphs[0] if template_id == 't2' else doc.add_paragraph()
                r = p.add_run() if template_id == 't2' else doc.paragraphs[0].add_run()
                r.add_picture(img_stream, width=Inches(1.0))
            
            # Left pane: Contact, Skills, Education
            left_cell.add_paragraph(cv_data.get("name", "Generated CV"), style="Title")
            if cv_data.get("email"): left_cell.add_paragraph(cv_data.get("email"))
            if cv_data.get("phone"): left_cell.add_paragraph(cv_data.get("phone"))
            add_skills(left_cell)
            add_edu(left_cell)

            # Right pane: Summary, Experience
            add_summary(right_cell)
            add_exp(right_cell)

        else: # Classic Top-down flow
            if img_stream:
                p = doc.add_paragraph()
                if align_center: p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.add_run().add_picture(img_stream, width=Inches(1.0))
            
            h = doc.add_heading(cv_data.get("name", "Generated CV"), 0)
            h.style.font.color.rgb = main_color
            if align_center: h.alignment = WD_ALIGN_PARAGRAPH.CENTER

            contact_info = [c for c in [cv_data.get("email"), cv_data.get("phone")] if c]
            if contact_info:
                p = doc.add_paragraph(" | ".join(contact_info))
                if align_center: p.alignment = WD_ALIGN_PARAGRAPH.CENTER

            add_summary(doc)
            add_exp(doc)
            add_edu(doc)
            add_skills(doc)
        
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        return StreamingResponse(
            file_stream, 
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
            headers={"Content-Disposition": "attachment; filename=resume.docx"}
        )
    
    return {"error": "Unsupported format from backend (PDF generation handled in frontend print)."}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
